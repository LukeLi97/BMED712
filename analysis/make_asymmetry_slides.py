"""Generate BMED712 asymmetry analysis slide deck (python-pptx)."""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── colour palette (Ocean Executive) ─────────────────────────────────
NAVY   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
MINT   = RGBColor(0x02, 0xC3, 0x9A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF4, 0xF8, 0xFB)
DARK   = RGBColor(0x0D, 0x1B, 0x2A)
GREY   = RGBColor(0x64, 0x74, 0x8B)
LIGHT_TEAL = RGBColor(0xCA, 0xE9, 0xF5)

# Group colors
C_HEALTHY = RGBColor(0x2E, 0x86, 0x3A)  # green
C_ORTHO   = RGBColor(0xE6, 0x8A, 0x00)  # amber
C_NEURO   = RGBColor(0xC0, 0x39, 0x2B)  # red

FIG_DIR = Path(__file__).resolve().parents[1] / "results" / "figures"
OUT     = Path(__file__).resolve().parents[1] / "results" / "asymmetry_analysis_slides.pptx"

W = Inches(13.33)
H = Inches(7.5)


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color: RGBColor, line_color=None, line_width=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = _Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT,
        font_face="Calibri", wrap=True, margin=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if margin is not None:
        txBox.text_frame.margin_left  = Inches(margin)
        txBox.text_frame.margin_right = Inches(margin)
        txBox.text_frame.margin_top   = Inches(margin)
        txBox.text_frame.margin_bottom= Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_face
    return txBox


def img(slide, path, x, y, w, h=None):
    path = str(path)
    if h is None:
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


def divider(slide, y=0.85, color=MINT, thickness=0.04):
    rect(slide, 0, y, 13.33, thickness, color)


# ── SLIDE 1: Title ───────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    bg(sl, DARK)

    # Left accent bar
    rect(sl, 0, 0, 0.18, 7.5, MINT)

    # Title
    txt(sl, "Gait Temporal Asymmetry Analysis",
        0.45, 1.5, 10.5, 1.4,
        size=44, bold=True, color=WHITE, font_face="Calibri")

    # Subtitle
    txt(sl, "Methods · Group Comparisons · Clinical Relevance",
        0.45, 3.0, 10.5, 0.7,
        size=22, italic=True, color=LIGHT_TEAL, font_face="Calibri")

    # Detail line
    txt(sl, "BMED712 Track A  |  216 subjects  |  974 trials  |  IMU Gait Dataset",
        0.45, 3.9, 10.5, 0.5,
        size=15, color=GREY, font_face="Calibri")

    # Bottom group pills
    for i, (label, c) in enumerate([("Healthy (n=70)", C_HEALTHY),
                                      ("Ortho (n=35)", C_ORTHO),
                                      ("Neuro (n=82)", C_NEURO)]):
        x = 0.45 + i * 2.6
        rect(sl, x, 5.8, 2.3, 0.6, c)
        txt(sl, label, x + 0.1, 5.85, 2.1, 0.5,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0)


# ── SLIDE 2: Background & Data Source ────────────────────────────────
def slide_background(prs):
    sl = blank_slide(prs)
    bg(sl, OFFWHITE)

    # Header bar
    rect(sl, 0, 0, 13.33, 1.05, NAVY)
    txt(sl, "Background & Data Source",
        0.35, 0.1, 12.0, 0.8,
        size=30, bold=True, color=WHITE, font_face="Calibri", margin=0)

    # Left column: context
    rect(sl, 0.3, 1.2, 5.8, 5.6, WHITE)
    txt(sl, "Clinical Context",
        0.5, 1.35, 5.4, 0.45,
        size=17, bold=True, color=NAVY, font_face="Calibri", margin=0)

    bullets_ctx = [
        "Gait asymmetry = difference in timing between left and right steps",
        "Healthy walking has natural lateralization (preferred side leads)",
        "Neurological & orthopedic conditions disrupt this pattern",
        "Heel-Strike (HS) timing from wearable IMUs → objective measure",
        "Key question: does |AI| differ across Healthy / Ortho / Neuro groups?",
    ]
    y = 1.95
    for b in bullets_ctx:
        rect(sl, 0.55, y + 0.05, 0.06, 0.06, MINT)
        txt(sl, b, 0.72, y, 5.2, 0.5, size=13, color=DARK, font_face="Calibri")
        y += 0.5

    # Right column: dataset
    rect(sl, 6.7, 1.2, 6.0, 5.6, WHITE)
    txt(sl, "Dataset",
        6.9, 1.35, 5.6, 0.45,
        size=17, bold=True, color=NAVY, font_face="Calibri", margin=0)

    # Stats grid
    stats = [
        ("216", "Subjects"),
        ("974", "Valid Trials"),
        ("4", "IMU Sensors"),
        ("100 Hz", "Sampling Rate"),
    ]
    for i, (val, label) in enumerate(stats):
        col = i % 2
        row = i // 2
        bx = 7.0 + col * 2.8
        by = 2.0 + row * 1.5
        rect(sl, bx, by, 2.4, 1.2, LIGHT_TEAL)
        txt(sl, val,   bx + 0.15, by + 0.05, 2.1, 0.7,
            size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER, margin=0)
        txt(sl, label, bx + 0.15, by + 0.72, 2.1, 0.4,
            size=12, color=GREY, align=PP_ALIGN.CENTER, margin=0)

    # Groups breakdown
    txt(sl, "Groups",
        6.9, 5.2, 5.6, 0.4,
        size=14, bold=True, color=NAVY, font_face="Calibri", margin=0)
    for i, (g, n, sub, c) in enumerate([
        ("Healthy", 70, "Control group", C_HEALTHY),
        ("Ortho",   35, "ACL, HOA, KOA", C_ORTHO),
        ("Neuro",   82, "PD, CVA, RIL, CIPN", C_NEURO),
    ]):
        bx = 7.0 + i * 1.95
        rect(sl, bx, 5.65, 1.75, 0.95, c)
        txt(sl, g,  bx + 0.08, 5.68, 1.6, 0.35,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0)
        txt(sl, f"n={n}", bx + 0.08, 5.98, 1.6, 0.28,
            size=11, color=WHITE, align=PP_ALIGN.CENTER, margin=0)


# ── SLIDE 3: Methods ─────────────────────────────────────────────────
def slide_methods(prs):
    sl = blank_slide(prs)
    bg(sl, OFFWHITE)

    rect(sl, 0, 0, 13.33, 1.05, TEAL)
    txt(sl, "Methods: How Asymmetry is Calculated",
        0.35, 0.1, 12.0, 0.8,
        size=30, bold=True, color=WHITE, font_face="Calibri", margin=0)

    # Step 1
    rect(sl, 0.3, 1.2, 0.6, 0.6, NAVY)
    txt(sl, "1", 0.3, 1.22, 0.6, 0.55, size=24, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, margin=0)
    txt(sl, "Extract Heel-Strike Events",
        1.05, 1.22, 7.0, 0.45, size=16, bold=True, color=NAVY, margin=0)
    txt(sl, "From metadata: leftGaitEvents / rightGaitEvents → pair[1] = HS sample index\n"
            "U-turn segment excluded via uturnBoundaries. Quality filter: ≥3 HS per side, no duplicates.",
        1.05, 1.65, 9.5, 0.7, size=12, color=DARK, margin=0)

    # Step 2
    rect(sl, 0.3, 2.6, 0.6, 0.6, NAVY)
    txt(sl, "2", 0.3, 2.62, 0.6, 0.55, size=24, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, margin=0)
    txt(sl, "Compute Step & Stride Times",
        1.05, 2.62, 7.0, 0.45, size=16, bold=True, color=NAVY, margin=0)
    txt(sl, "Step time = interval R→next L (left step) and L→next R (right step)\n"
            "Stride time = consecutive same-foot HS intervals. Valid range: 0.2–3.0 s step, 0.4–5.0 s stride.",
        1.05, 3.05, 9.5, 0.7, size=12, color=DARK, margin=0)

    # Step 3 — metrics table
    rect(sl, 0.3, 4.0, 0.6, 0.6, NAVY)
    txt(sl, "3", 0.3, 4.02, 0.6, 0.55, size=24, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, margin=0)
    txt(sl, "Five Asymmetry Metrics",
        1.05, 4.02, 7.0, 0.45, size=16, bold=True, color=NAVY, margin=0)

    # Metrics table
    headers = ["Metric", "Formula", "Interpretation"]
    rows_data = [
        ["Signed AI",   "(L−R) / mean(L,R)",         "Direction of asymmetry (positive = L dominates)"],
        ["|AI|",        "|L−R| / mean(L,R)",          "Magnitude — our PRIMARY metric"],
        ["|L−R| diff",  "|L−R| (seconds)",             "Absolute timing gap"],
        ["Ratio",       "L / R",                       "Ratio of left to right duration"],
        ["CV",          "SD / mean (per side)",        "Within-trial variability"],
    ]
    col_w = [1.6, 3.0, 4.8]
    col_x = [1.05, 2.65, 5.65]
    row_h = 0.45
    start_y = 4.55

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        rect(sl, cx, start_y, cw - 0.05, row_h, NAVY)
        txt(sl, hdr, cx + 0.08, start_y + 0.05, cw - 0.18, row_h - 0.1,
            size=12, bold=True, color=WHITE, margin=0)

    for i, row in enumerate(rows_data):
        by = start_y + row_h + i * row_h
        fill = LIGHT_TEAL if i % 2 == 0 else WHITE
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            rect(sl, cx, by, cw - 0.05, row_h - 0.03, fill)
            bold = (j == 0)
            c_color = NAVY if j == 1 else DARK
            txt(sl, cell, cx + 0.08, by + 0.05, cw - 0.18, row_h - 0.12,
                size=11, bold=bold, color=c_color, margin=0)

    # Subject-level note
    rect(sl, 0.3, 7.05, 12.7, 0.38, MINT)
    txt(sl, "⚠  Subject-level aggregation: trial means averaged per subject → 216 data points (avoids pseudoreplication)",
        0.45, 7.07, 12.4, 0.32, size=12, bold=True, color=DARK, margin=0)


# ── SLIDE 4: Key Finding ─────────────────────────────────────────────
def slide_key_finding(prs):
    sl = blank_slide(prs)
    bg(sl, DARK)
    rect(sl, 0, 0, 0.18, 7.5, MINT)

    txt(sl, "Key Finding: Healthy Gait is MORE Asymmetric",
        0.45, 0.25, 12.5, 0.9,
        size=36, bold=True, color=WHITE, font_face="Calibri", margin=0)

    txt(sl, "Counterintuitive but validated: neurological pathology leads to loss of natural motor lateralization → more symmetric gait",
        0.45, 1.15, 10.5, 0.5,
        size=14, italic=True, color=LIGHT_TEAL, font_face="Calibri", margin=0)

    # Main boxplot figure
    fig = FIG_DIR / "step06_asymmetry_boxplot_absAI.png"
    if fig.exists():
        img(sl, fig, 0.45, 1.75, 6.5, 5.0)

    # Stats callout boxes on the right
    callouts = [
        ("Cohen's d = 0.77", "H vs N (Stride |AI|)", MINT),
        ("p < 0.001", "Kruskal-Wallis", TEAL),
        ("AUC = 0.716", "Single-feature screening", RGBColor(0x1C, 0x72, 0x93)),
        ("83%", "Specificity at optimal threshold", RGBColor(0x02, 0x8A, 0x6F)),
    ]
    for i, (val, label, c) in enumerate(callouts):
        by = 1.75 + i * 1.3
        rect(sl, 7.3, by, 5.7, 1.15, c)
        txt(sl, val,   7.5, by + 0.05, 5.3, 0.65,
            size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0)
        txt(sl, label, 7.5, by + 0.7,  5.3, 0.38,
            size=13, color=WHITE, align=PP_ALIGN.CENTER, margin=0)


# ── SLIDE 5: Statistical Validation ──────────────────────────────────
def slide_stats(prs):
    sl = blank_slide(prs)
    bg(sl, OFFWHITE)

    rect(sl, 0, 0, 13.33, 1.05, NAVY)
    txt(sl, "Statistical Validation",
        0.35, 0.1, 12.0, 0.8,
        size=30, bold=True, color=WHITE, font_face="Calibri", margin=0)

    # Bootstrap LME figure
    fig = FIG_DIR / "step11_bootstrap_lme.png"
    if fig.exists():
        img(sl, fig, 0.3, 1.15, 8.2, 5.8)

    # Results column
    results = [
        ("Bootstrap CI (n=10,000)", "d = 0.77  [0.50, 1.07]", "AUC = 0.716  [0.635, 0.792]"),
        ("Linear Mixed Effects Model", "stride_AI ~ group + (1|Subject)", "Neuro: β=−0.019, p=3×10⁻⁶"),
        ("Neuro vs Healthy", "Significant loss of lateralization", "Ortho: β≈0, p=1.00 (not affected)"),
    ]
    for i, (header, line1, line2) in enumerate(results):
        by = 1.3 + i * 1.85
        rect(sl, 8.75, by, 4.3, 1.6, WHITE)
        rect(sl, 8.75, by, 0.1, 1.6, TEAL)
        txt(sl, header, 9.0, by + 0.1, 3.9, 0.4,
            size=13, bold=True, color=NAVY, margin=0)
        txt(sl, line1,  9.0, by + 0.5, 3.9, 0.38,
            size=12, color=DARK, margin=0)
        txt(sl, line2,  9.0, by + 0.88, 3.9, 0.38,
            size=12, italic=True, color=TEAL, margin=0)

    # Subtype bar note
    rect(sl, 8.75, 6.8, 4.3, 0.6, LIGHT_TEAL)
    txt(sl, "Best subtype:  RIL d=0.87  >  PD d=0.77  >  CVA d=0.73",
        8.85, 6.85, 4.1, 0.5, size=12, bold=True, color=NAVY, margin=0)


# ── SLIDE 6: Clinical Relevance ───────────────────────────────────────
def slide_clinical(prs):
    sl = blank_slide(prs)
    bg(sl, OFFWHITE)

    rect(sl, 0, 0, 13.33, 1.05, TEAL)
    txt(sl, "Clinical Relevance",
        0.35, 0.1, 12.0, 0.8,
        size=30, bold=True, color=WHITE, font_face="Calibri", margin=0)

    # VGA scatter
    fig_vga = FIG_DIR / "step13_vga_stride_absAI.png"
    if fig_vga.exists():
        img(sl, fig_vga, 0.3, 1.15, 6.4, 4.5)

    # Phase figure
    fig_phase = FIG_DIR / "step12_phase_stride_absAI.png"
    if fig_phase.exists():
        img(sl, fig_phase, 6.9, 1.15, 6.1, 4.5)

    # Bottom bar
    rect(sl, 0, 5.9, 13.33, 1.55, NAVY)
    findings = [
        ("VGA Correlation", "r = −0.206,  p = 2.3×10⁻¹⁰\n(n=927 trials)"),
        ("Phase Analysis", "Pre-uturn p=0.026\nU-turn: ns — no added value"),
        ("Ortho U-turn Spike", "|AI| = 0.110 vs 0.005 (healthy)\nCompensatory weight-shift"),
    ]
    for i, (header, detail) in enumerate(findings):
        bx = 0.3 + i * 4.35
        rect(sl, bx, 6.0, 4.1, 1.3, TEAL)
        txt(sl, header, bx + 0.12, 6.05, 3.86, 0.4,
            size=13, bold=True, color=WHITE, margin=0)
        txt(sl, detail, bx + 0.12, 6.45, 3.86, 0.75,
            size=12, color=LIGHT_TEAL, margin=0)


# ── SLIDE 7: Conclusions ──────────────────────────────────────────────
def slide_conclusions(prs):
    sl = blank_slide(prs)
    bg(sl, DARK)
    rect(sl, 0, 0, 0.18, 7.5, MINT)

    txt(sl, "Conclusions",
        0.45, 0.3, 12.5, 0.8,
        size=40, bold=True, color=WHITE, font_face="Calibri", margin=0)

    points = [
        ("Healthy > Pathological asymmetry",
         "Natural motor lateralization is lost in neurological disease — not gained"),
        ("Stride |AI| is the strongest discriminator",
         "Cohen's d = 0.77 [0.50, 1.07], AUC = 0.716 [0.635, 0.792] (bootstrap confirmed)"),
        ("Neurological ≠ Orthopedic",
         "LME: neuro β=−0.019 (p=3×10⁻⁶), ortho β≈0 (p=1.00) — different mechanisms"),
        ("Clinically validated",
         "Spearman r=−0.206 with VGA clinician rating (p=2.3×10⁻¹⁰, n=927 trials)"),
        ("Steady-state walking carries the signal",
         "Pre-uturn phase: p=0.026; U-turn adds noise not information"),
    ]
    for i, (heading, detail) in enumerate(points):
        by = 1.25 + i * 1.15
        rect(sl, 0.45, by, 12.5, 1.0, RGBColor(0x0D, 0x2D, 0x45))
        rect(sl, 0.45, by, 0.12, 1.0, MINT)
        txt(sl, heading, 0.7, by + 0.04, 12.0, 0.38,
            size=15, bold=True, color=WHITE, margin=0)
        txt(sl, detail,  0.7, by + 0.42, 12.0, 0.45,
            size=13, italic=False, color=LIGHT_TEAL, margin=0)


# ── MAIN ──────────────────────────────────────────────────────────────
def main():
    prs = prs_new()
    slide_title(prs)
    slide_background(prs)
    slide_methods(prs)
    slide_key_finding(prs)
    slide_stats(prs)
    slide_clinical(prs)
    slide_conclusions(prs)
    prs.save(str(OUT))
    print(f"Saved → {OUT}")


if __name__ == "__main__":
    main()
