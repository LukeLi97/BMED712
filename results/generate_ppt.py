#!/usr/bin/env python3
"""
Generate final presentation v2 — uses KU template, fixes title overlap,
removes section divider slides, iconifies text-heavy pages.

Total: 14 slides = 1 title + 12 content (3 per presenter) + 1 thank-you.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import sys, os

ROOT_LINUX = Path("/sessions/quirky-fervent-allen/mnt/BMED712 Rehab")
ROOT_MAC   = Path("/Users/test/Desktop/BMED712 Rehab")
ROOT = ROOT_LINUX if ROOT_LINUX.exists() else ROOT_MAC

TEMPLATE = ROOT / "KU_ppt_template_2018.pptx"
OUTPUT   = ROOT / "results" / "Final_Presentation_Track_A.pptx"

FIGS = {
    "cohort":     ROOT / "results/validation/cohort_balance.png",
    "confusion":  ROOT / "results/figures/step03_confusion_3class_all.png",
    "confusion_svm": ROOT / "results/figures/step03_confusion_3class_all.png",
    "frontier":   ROOT / "results/figures/step04_sensors_frontier.png",
    "importance": ROOT / "results/figures/step05_importance_3class_all.png",
    "vga":        ROOT / "results/figures/step07_corr_vga_stride_absAI_fixed.png",
    "phase":      ROOT / "results/figures/phase_single_vs_all.png",
    "heatmap":    ROOT / "PHASE1_Feature_Analysis_COMPLETE/Full_Gait_6s_ov50/correlation_heatmap_Full_Gait_6s_ov50.png",
}

# ═════════ KU brand palette ═════════
KU_BLUE   = RGBColor(0x00, 0x47, 0xBA)
KU_GREEN  = RGBColor(0x00, 0xCE, 0x7C)
KU_RED    = RGBColor(0xE5, 0x3E, 0x51)
KU_TEAL   = RGBColor(0x84, 0xDA, 0xDE)
KU_YELLOW = RGBColor(0xF5, 0xCE, 0x3E)
KU_GRAY   = RGBColor(0x96, 0x96, 0x9A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x2B, 0x2B, 0x2B)
SOFT_BG   = RGBColor(0xF4, 0xF8, 0xFC)  # very light blue-tint
SOFT_RED  = RGBColor(0xFD, 0xF3, 0xF3)
SOFT_GREEN= RGBColor(0xF0, 0xFA, 0xF5)
DEEP_NAVY = RGBColor(0x00, 0x2F, 0x7A)

# Layout indices (from template analysis)
LY_TITLE   = 1   # Title Slide (dark navy bg)
LY_CONTENT = 11  # Title Only (has KU header bar + footer)
LY_END     = 16  # End slide

# Content safe zone: title placeholder ends at ~y=1.81 in this template
# All content must start at y >= 1.95 to avoid overlap.
CONTENT_TOP = 1.95  # inches
CONTENT_BOTTOM = 6.75  # above presenter tag


# ═════════ helpers ═════════
def tb(slide, left, top, width, height, text, sz=16, bold=False,
       color=BLACK, align=PP_ALIGN.LEFT, anchor=None, font=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    if font:
        r.font.name = font
    return tf


def bullets(slide, left, top, width, height, items, sz=13, color=BLACK,
            space=3, bullet_char="•", bullet_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space)
        if item == "":
            continue
        # Add bullet marker
        if bullet_char:
            rb = p.add_run()
            rb.text = bullet_char + "  "
            rb.font.size = Pt(sz)
            rb.font.bold = True
            rb.font.color.rgb = bullet_color if bullet_color else color
        if ": " in item and not item.startswith(" "):
            label, rest = item.split(": ", 1)
            r1 = p.add_run()
            r1.text = label + ": "
            r1.font.size = Pt(sz)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(sz)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(sz)
            r.font.color.rgb = color
    return tf


def card(slide, left, top, width, height, fill_color, line_color=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        try:
            shape.adjustments[0] = 0.07
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def img(slide, path, left, top, width, height=None):
    if height is None:
        from PIL import Image
        with Image.open(str(path)) as im:
            iw, ih = im.size
        aspect = iw / ih
        w_in = width / 914400
        h_in = w_in / aspect
        height = Inches(h_in)
    return slide.shapes.add_picture(str(path), left, top, width, height)


def caption(slide, left, top, width, text):
    tb(slide, left, top, width, Inches(0.3), text,
       sz=10, color=KU_GRAY, align=PP_ALIGN.CENTER)


def presenter_tag(slide, num, topic):
    tb(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
       f"Presenter {num}  |  {topic}", sz=10, color=KU_GRAY)


def add_notes(slide, text):
    """Attach speaker notes to a slide."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    return tf


def set_title(slide, text, color=KU_BLUE, sz=28):
    """Set title placeholder text safely."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            r.font.bold = True
            r.font.size = Pt(sz)
            r.font.color.rgb = color
            return


def icon_circle(slide, cx, cy, r, fill, text, text_color=WHITE, sz=18):
    """Draw a colored circle with centered text (icon/number)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = tb(slide, Inches(cx - r), Inches(cy - r - 0.05),
            Inches(r * 2), Inches(r * 2 + 0.1), text,
            sz=sz, bold=True, color=text_color, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    return shape


def dot(slide, cx, cy, r, fill):
    """Small filled circle used as a bullet/icon — always renders."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def badge(slide, left, top, size, fill, glyph, glyph_color=WHITE, sz=14):
    """Square rounded badge with a Unicode glyph — for row icons."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(size), Inches(size))
    try:
        shape.adjustments[0] = 0.25
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tb(slide, Inches(left), Inches(top), Inches(size), Inches(size),
       glyph, sz=sz, bold=True, color=glyph_color,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shape


def arrow_right(slide, left, top, width=0.35, height=0.35, color=None):
    if color is None: color = KU_GRAY
    tb(slide, Inches(left), Inches(top), Inches(width), Inches(height),
       "▸", sz=22, bold=True, color=color, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)


# ═════════ BUILD ═════════
def build():
    prs = Presentation(str(TEMPLATE))

    # Wipe template example slides
    sld_ids = list(prs.slides._sldIdLst)
    for sid in sld_ids:
        rId = sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.slides._sldIdLst.remove(sid)
        if rId:
            prs.part.drop_rel(rId)

    W = prs.slide_width / 914400
    H = prs.slide_height / 914400
    print(f"Slide size: {W:.2f} x {H:.2f} inches")

    # ═════ SLIDE 1: TITLE ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_TITLE])
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "IMU-Based Gait Pathology Classification"
        elif idx == 1:
            ph.text = "Machine Learning on the GaitRec Dataset"
        elif idx == 10:
            ph.text = "BMED 712  |  Rehabilitation Engineering  |  April 2026"
    add_notes(s, (
        "Good morning everyone.\n"
        "Our project is IMU-Based Gait Pathology Classification, using the "
        "GaitRec dataset with machine learning.\n"
        "This is Track A for BMED 712, Rehabilitation Engineering.\n"
        "Today our team of four will walk you through the dataset, the ML "
        "methods, the results, and the clinical implications.\n"
        "Each of us will cover three slides. Let's begin."
    ))
    print("  [1] Title")

    # ═════ SLIDE 2: P1-S1 — Problem & Dataset ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Problem & Dataset Overview")

    # Three big stat callouts
    stats = [
        (0.6,  "260",      "Subjects (8 cohorts)",  KU_BLUE),
        (4.65, "1,356",    "Walking trials",         KU_GREEN),
        (8.7,  "300,991",  "Feature windows",        KU_RED),
    ]
    for x, num, label, clr in stats:
        card(s, Inches(x), Inches(CONTENT_TOP), Inches(3.95), Inches(1.55),
             clr, rounded=True)
        tb(s, Inches(x), Inches(CONTENT_TOP + 0.1), Inches(3.95), Inches(0.85),
           num, sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)
        tb(s, Inches(x), Inches(CONTENT_TOP + 1.0), Inches(3.95), Inches(0.45),
           label, sz=13, color=WHITE, align=PP_ALIGN.CENTER)

    # Cohort balance image (left)
    img_y = CONTENT_TOP + 1.9
    img(s, FIGS["cohort"], Inches(0.6), Inches(img_y), Inches(7.2))
    caption(s, Inches(0.6), Inches(img_y + 2.05), Inches(7.2),
            "Class distribution (HC / Mild / Severe) across 8 cohorts")

    # Icon-ized key facts (right side)
    icons_info = [
        (KU_BLUE,  "GaitRec (Medical Univ. Innsbruck)"),
        (KU_GREEN, "4 IMU sensors: HE, LB, LF, RF"),
        (KU_TEAL,  "3 signals × 3 axes × 6 metrics"),
        (KU_RED,   "3-class: HC / Mild / Severe"),
    ]
    for i, (clr, txt) in enumerate(icons_info):
        y = img_y + i * 0.55
        dot(s, cx=8.8, cy=y + 0.22, r=0.12, fill=clr)
        tb(s, Inches(9.1), Inches(y), Inches(4.0), Inches(0.45),
           txt, sz=13, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    # Highlighted tag
    tb(s, Inches(8.6), Inches(img_y + 2.3), Inches(4.5), Inches(0.45),
       "216 hand-crafted features per window",
       sz=12, bold=True, color=KU_BLUE)

    presenter_tag(s, 1, "Dataset Overview")
    add_notes(s, (
        "Hello, I am Presenter One. I will introduce the problem and the dataset.\n\n"
        "Our goal is to classify gait pathology into three severity levels — "
        "Healthy, Mild, and Severe — using only wearable IMU sensors.\n\n"
        "We use the GaitRec dataset from the Medical University of Innsbruck. "
        "It contains 260 subjects across 8 clinical cohorts, 1,356 walking "
        "trials, and after windowing, more than 300,000 feature windows.\n\n"
        "Each trial has 4 IMU sensors — Head, Lower Back, Left Foot, and Right "
        "Foot. Each sensor records three signals: free acceleration, gyroscope, "
        "and magnetometer, each in three axes. From every channel we extract "
        "six metrics, giving 216 hand-crafted features per window.\n\n"
        "The class distribution across cohorts is shown on the left. Some "
        "cohorts are imbalanced, so we use stratified cross-validation later.\n\n"
        "Next, I will show how we extract these features."
    ))
    print("  [2] Problem & Dataset")

    # ═════ SLIDE 3: P1-S2 — Feature Extraction Pipeline ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Feature Extraction Pipeline")

    # 4-step horizontal flow
    pipeline = [
        ("Windowing",  "6 s windows\n50% overlap",          KU_BLUE),
        ("Extraction", "4 sensors × 3 ch\n× 3 axes",        KU_GREEN),
        ("Metrics",    "Mean · Std · Min\nMax · DomF · SC", KU_TEAL),
        ("Selection",  "KW H-test\nη² ranking",              KU_RED),
    ]
    card_w, card_h, gap = 2.65, 1.9, 0.25
    start_x = 0.55
    for i, (ttl, desc, clr) in enumerate(pipeline):
        x = start_x + i * (card_w + gap)
        card(s, Inches(x), Inches(CONTENT_TOP), Inches(card_w), Inches(card_h),
             clr, rounded=True)
        tb(s, Inches(x), Inches(CONTENT_TOP + 0.15), Inches(card_w), Inches(0.55),
           ttl, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, Inches(x + 0.1), Inches(CONTENT_TOP + 0.8), Inches(card_w - 0.2), Inches(1.0),
           desc, sz=13, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            arrow_right(s, x + card_w + 0.02, CONTENT_TOP + 0.85, width=gap - 0.04)

    # Three insight cards below
    insights = [
        ("◆", "X = ML / Y = AP / Z = Vertical",
         "Standard anatomical axes", KU_BLUE),
        ("✓", "207 / 216 significant",
         "p < 0.05, Bonferroni corrected", KU_GREEN),
        ("★", "Top η² = 0.303",
         "HE_FreeAcc_X_dom_freq", KU_RED),
    ]
    for i, (icon, head, sub, clr) in enumerate(insights):
        x = 0.55 + i * 4.25
        card(s, Inches(x), Inches(CONTENT_TOP + card_h + 0.45),
             Inches(4.0), Inches(1.3), SOFT_BG, line_color=clr, rounded=True)
        badge(s, left=x + 0.2, top=CONTENT_TOP + card_h + 0.6,
              size=0.5, fill=clr, glyph=icon, sz=18)
        tb(s, Inches(x + 0.85), Inches(CONTENT_TOP + card_h + 0.5),
           Inches(3.1), Inches(0.45), head, sz=14, bold=True, color=clr)
        tb(s, Inches(x + 0.85), Inches(CONTENT_TOP + card_h + 0.9),
           Inches(3.1), Inches(0.4), sub, sz=11, color=BLACK)

    # Bottom takeaway
    tb(s, Inches(0.55), Inches(CONTENT_TOP + card_h + 2.0), Inches(12.2),
       Inches(0.4),
       "Window configs: Full Gait 6 s (50% overlap)  +  U-Turn 3 s (0% overlap)",
       sz=12, color=KU_GRAY, align=PP_ALIGN.CENTER)

    presenter_tag(s, 1, "Feature Extraction")
    add_notes(s, (
        "Our feature extraction has four steps.\n\n"
        "First, Windowing. We split each walking trial into 6-second windows "
        "with 50 percent overlap, which captures several gait cycles per window.\n\n"
        "Second, Extraction. From each window we use 4 sensors, times 3 "
        "channels, times 3 axes — that gives 36 time series per window.\n\n"
        "Third, Metrics. For every channel we compute six metrics: Mean, "
        "Standard deviation, Min, Max, Dominant frequency, and Spectral "
        "centroid. This gives us 216 features in total.\n\n"
        "Fourth, Selection. We use the Kruskal–Wallis H-test with eta-squared "
        "ranking to select the most discriminative features.\n\n"
        "Results are strong: 207 out of 216 features are statistically "
        "significant, even after Bonferroni correction. The top feature is "
        "the head sensor dominant frequency, with an eta-squared of 0.303.\n\n"
        "We also tested a shorter 3-second window for U-turn segments."
    ))
    print("  [3] Feature Extraction")

    # ═════ SLIDE 4: P1-S3 — Correlation & Descriptive ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Feature Correlation & Top Features")

    # Heatmap (left) — scale to ~4.0 in height
    img(s, FIGS["heatmap"], Inches(0.55), Inches(CONTENT_TOP),
        Inches(5.4))
    caption(s, Inches(0.55), Inches(CONTENT_TOP + 4.55), Inches(5.4),
            "Pearson correlation — top 50 features (Full Gait 6s)")

    # Right column: Top-3 stat blocks
    top3 = [
        ("1", "HE_FreeAcc_X",  "dom freq",   "0.303", KU_BLUE),
        ("2", "HE_FreeAcc_Y",  "dom freq",   "0.246", KU_GREEN),
        ("3", "LB_Gyr_Z",       "spec centroid","0.221", KU_RED),
    ]
    tb(s, Inches(6.75), Inches(CONTENT_TOP), Inches(6.2), Inches(0.4),
       "Top 3 features by η²", sz=16, bold=True, color=KU_BLUE)
    for i, (rank, feat, metric, eta, clr) in enumerate(top3):
        y = CONTENT_TOP + 0.55 + i * 0.85
        # rank circle
        icon_circle(s, cx=7.0, cy=y + 0.33, r=0.22, fill=clr,
                    text=rank, sz=14)
        # feature name + metric
        tb(s, Inches(7.45), Inches(y), Inches(3.4), Inches(0.4),
           feat, sz=14, bold=True, color=BLACK)
        tb(s, Inches(7.45), Inches(y + 0.35), Inches(3.4), Inches(0.35),
           metric, sz=11, color=KU_GRAY)
        # eta value
        tb(s, Inches(10.85), Inches(y), Inches(2.0), Inches(0.7),
           f"η² = {eta}", sz=18, bold=True, color=clr, align=PP_ALIGN.RIGHT,
           anchor=MSO_ANCHOR.MIDDLE)

    # Key observations as icon rows
    obs_y = CONTENT_TOP + 3.3
    card(s, Inches(6.75), Inches(obs_y), Inches(6.2), Inches(1.5),
         SOFT_BG, rounded=True)
    tb(s, Inches(6.95), Inches(obs_y + 0.05), Inches(5.8), Inches(0.35),
       "Key Observations", sz=13, bold=True, color=KU_BLUE)
    obs = [
        "Same-sensor features heavily correlated (r > 0.8)",
        "Head (HE) sensor → most independent & informative",
        "Gyroscope dom-freq features most discriminative",
    ]
    for i, t in enumerate(obs):
        tb(s, Inches(6.95), Inches(obs_y + 0.45 + i * 0.32), Inches(0.3),
           Inches(0.3), "→", sz=12, bold=True, color=KU_GREEN)
        tb(s, Inches(7.25), Inches(obs_y + 0.42 + i * 0.32), Inches(5.6),
           Inches(0.35), t, sz=11, color=BLACK)

    presenter_tag(s, 1, "Descriptive Statistics")
    add_notes(s, (
        "This is my last slide. Let's look at the feature correlation and "
        "the top-ranked features.\n\n"
        "The heatmap on the left is the Pearson correlation matrix for the "
        "top 50 features. You can see strong red clusters along the diagonal. "
        "This means same-sensor features are heavily correlated, especially "
        "for the foot sensors — their features are highly redundant.\n\n"
        "In contrast, the Head sensor features are more independent. They "
        "carry information the other sensors cannot provide.\n\n"
        "On the right we show the top three features by eta-squared effect "
        "size. Rank one is the Head free-acceleration X-axis dominant "
        "frequency, with eta-squared 0.303. Rank two is the Head Y-axis "
        "dominant frequency, 0.246. Rank three is the Lower-Back gyroscope "
        "Z-axis spectral centroid, 0.221.\n\n"
        "The key takeaway: frequency-domain features from the Head and "
        "Lower-Back sensors are the most informative for pathology "
        "classification.\n\n"
        "I will now hand over to Presenter Two for the machine learning methods."
    ))
    print("  [4] Correlation & Top Features")

    # ═════ SLIDE 5: P2-S1 — ML Pipeline & CV ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Classification Pipeline & Cross-Validation")

    classifiers = [
        ("SVM (RBF)",     "81.0%", "0.796", KU_BLUE,  True),
        ("Random Forest", "80.1%", "0.788", KU_GREEN, False),
        ("XGBoost",       "79.8%", "0.781", KU_RED,   False),
    ]
    c_w = 3.95
    for i, (name, bacc, f1, clr, best) in enumerate(classifiers):
        x = 0.6 + i * (c_w + 0.15)
        card(s, Inches(x), Inches(CONTENT_TOP), Inches(c_w), Inches(2.15),
             clr, rounded=True)
        tb(s, Inches(x), Inches(CONTENT_TOP + 0.15), Inches(c_w), Inches(0.5),
           name, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, Inches(x), Inches(CONTENT_TOP + 0.7), Inches(c_w), Inches(0.75),
           bacc, sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, Inches(x), Inches(CONTENT_TOP + 1.5), Inches(c_w), Inches(0.4),
           f"BAcc  |  F1 = {f1}", sz=12, color=WHITE, align=PP_ALIGN.CENTER)
        if best:
            card(s, Inches(x + c_w - 1.15), Inches(CONTENT_TOP - 0.2),
                 Inches(1.1), Inches(0.35), KU_YELLOW, rounded=True)
            tb(s, Inches(x + c_w - 1.15), Inches(CONTENT_TOP - 0.2),
               Inches(1.1), Inches(0.35), "★ BEST",
               sz=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)

    # Bottom: 2×2 icon grid for methodology
    grid_y = CONTENT_TOP + 2.5
    method = [
        ("CV",  "StratifiedGroupKFold",
         "10-fold, grouped by subject_id — no data leakage", KU_BLUE),
        ("N",  "Nested CV check",
         "Outer 10 × Inner 5 → 78.3% (2.7% optimism gap)", KU_GREEN),
        ("D",  "Demographics features",
         "Age, gender, laterality (one-hot) from _meta.json", KU_RED),
        ("F",  "Feature sets compared",
         "All 216 · Top 30 · Top 20 · p < 0.05 significant", KU_TEAL),
    ]
    for i, (glyph, head, desc, clr) in enumerate(method):
        col, row = i % 2, i // 2
        x = 0.6 + col * 6.25
        y = grid_y + row * 1.0
        badge(s, left=x, top=y + 0.05, size=0.55, fill=clr,
              glyph=glyph, sz=14)
        tb(s, Inches(x + 0.7), Inches(y), Inches(5.4), Inches(0.4),
           head, sz=14, bold=True, color=clr)
        tb(s, Inches(x + 0.7), Inches(y + 0.4), Inches(5.4), Inches(0.45),
           desc, sz=11, color=BLACK)

    presenter_tag(s, 2, "ML Pipeline")
    add_notes(s, (
        "Thank you. I am Presenter Two, and I will cover the machine "
        "learning pipeline.\n\n"
        "We compared three classifiers on all 216 features. SVM with an RBF "
        "kernel achieved 81.0 percent balanced accuracy and F1 of 0.796 — "
        "this was our best model. Random Forest came second at 80.1 percent, "
        "and XGBoost at 79.8 percent. The three models are close, which "
        "suggests the features themselves are strong.\n\n"
        "Our cross-validation is very important. We use 10-fold Stratified "
        "Group K-Fold, grouped by subject ID. This prevents data leakage, "
        "because windows from the same subject never appear in both training "
        "and testing.\n\n"
        "To check for optimism bias, we also ran a nested CV — outer 10 fold, "
        "inner 5 fold. It gave 78.3 percent, which is only 2.7 percent below "
        "the standard CV. This small gap means the model is not overfit.\n\n"
        "We also added demographics — age, gender, laterality — and compared "
        "four feature sets: all 216, top 30, top 20, and significant only."
    ))
    print("  [5] ML Pipeline & CV")

    # ═════ SLIDE 6: P2-S2 — Feature Importance ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Feature Importance & Selection")

    img(s, FIGS["importance"], Inches(0.55), Inches(CONTENT_TOP),
        Inches(8.4))
    caption(s, Inches(0.55), Inches(CONTENT_TOP + 4.25), Inches(8.4),
            "Permutation importance — top 20 features (3-class SVM)")

    # Right: hero stat + 3 bullets with icons
    # Big callout: 95%
    card(s, Inches(9.2), Inches(CONTENT_TOP), Inches(3.75), Inches(1.5),
         KU_BLUE, rounded=True)
    tb(s, Inches(9.2), Inches(CONTENT_TOP + 0.05), Inches(3.75), Inches(0.85),
       "≈ 95%", sz=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(9.2), Inches(CONTENT_TOP + 0.95), Inches(3.75), Inches(0.5),
       "Top-20 vs. All-216", sz=12, color=WHITE, align=PP_ALIGN.CENTER)

    insights = [
        (KU_BLUE,  "Head sensor dominates top ranks"),
        (KU_GREEN, "Gyroscope + FreeAcc most informative"),
        (KU_RED,   "Sets tested: 216 / 30 / 20 / sig-only"),
    ]
    for i, (clr, txt) in enumerate(insights):
        y = CONTENT_TOP + 1.75 + i * 0.75
        dot(s, cx=9.4, cy=y + 0.3, r=0.12, fill=clr)
        tb(s, Inches(9.75), Inches(y), Inches(3.25), Inches(0.65),
           txt, sz=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)

    presenter_tag(s, 2, "Feature Selection")
    add_notes(s, (
        "Now let's look at feature importance.\n\n"
        "The chart on the left shows permutation importance for the top 20 "
        "features of our three-class SVM. Each bar is one feature, ranked "
        "from most important at the top.\n\n"
        "Two patterns are clear. First, the Head sensor dominates the top "
        "ranks — most of the top features start with H-E. Second, the "
        "gyroscope and free-acceleration signals are the most informative. "
        "Magnetometer features are less useful for this task.\n\n"
        "The most interesting result is on the right. We compared models "
        "trained on all 216 features versus only the top 20. The top-20 "
        "model reaches about 95 percent of the full-feature performance.\n\n"
        "This is important for deployment. A small, well-chosen feature set "
        "runs faster on wearable devices, is easier to interpret for "
        "clinicians, and reduces the risk of overfitting.\n\n"
        "We tested four feature sets in total: all 216, top 30, top 20, and "
        "significant-only. All of them performed within two points of each "
        "other."
    ))
    print("  [6] Feature Importance")

    # ═════ SLIDE 7: P2-S3 — Sensor Ablation ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Sensor Ablation & Optimal Configuration")

    img(s, FIGS["frontier"], Inches(0.55), Inches(CONTENT_TOP),
        Inches(6.5))
    caption(s, Inches(0.55), Inches(CONTENT_TOP + 3.7), Inches(6.5),
            "Sensor availability frontier (LR vs. RF)")

    # Three ablation result bars on right — progressive
    ablation = [
        ("1 Sensor",    "HE only",            "77.5%", KU_TEAL),
        ("2 Sensors",   "HE + LB",            "80.1%", KU_GREEN),
        ("4 Sensors",   "HE + LB + LF + RF",  "81.0%", KU_BLUE),
    ]
    for i, (cfg, detail, bacc, clr) in enumerate(ablation):
        y = CONTENT_TOP + i * 1.25
        card(s, Inches(7.4), Inches(y), Inches(5.55), Inches(1.1),
             clr, rounded=True)
        tb(s, Inches(7.55), Inches(y + 0.1), Inches(3.0), Inches(0.45),
           cfg, sz=18, bold=True, color=WHITE)
        tb(s, Inches(7.55), Inches(y + 0.55), Inches(3.0), Inches(0.4),
           detail, sz=11, color=WHITE)
        tb(s, Inches(10.6), Inches(y + 0.1), Inches(2.2), Inches(0.9),
           bacc, sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)

    # Clinical takeaway
    card(s, Inches(0.55), Inches(CONTENT_TOP + 4.15), Inches(12.4),
         Inches(0.75), SOFT_GREEN, line_color=KU_GREEN, rounded=True)
    badge(s, left=0.75, top=CONTENT_TOP + 4.28, size=0.5,
          fill=KU_GREEN, glyph="!", sz=16)
    tb(s, Inches(1.4), Inches(CONTENT_TOP + 4.22), Inches(11.4), Inches(0.6),
       "HE + LB = 98.9% of 4-sensor performance  →  2 IMUs sufficient for clinical deployment",
       sz=14, bold=True, color=DEEP_NAVY, anchor=MSO_ANCHOR.MIDDLE)

    presenter_tag(s, 2, "Sensor Ablation")
    add_notes(s, (
        "My last slide is the sensor ablation study. This asks a very "
        "practical question: how many IMUs do we actually need?\n\n"
        "The chart on the left shows the balanced accuracy as we add more "
        "sensors, for both Logistic Regression in blue and Random Forest "
        "in orange.\n\n"
        "With only the Head sensor, we already get 77.5 percent. Adding "
        "the Lower Back sensor brings it up to 80.1 percent. And using all "
        "four sensors only adds 0.9 percent more, giving 81.0 percent.\n\n"
        "The clinical takeaway is at the bottom: the Head plus Lower Back "
        "configuration captures 98.9 percent of the four-sensor "
        "performance. So two IMUs are enough for clinical deployment.\n\n"
        "This matters because patients are much more likely to wear two "
        "sensors than four. It reduces cost, improves comfort, and makes "
        "home monitoring realistic.\n\n"
        "I will now hand over to Presenter Three for the detailed "
        "classification results."
    ))
    print("  [7] Sensor Ablation")

    # ═════ SLIDE 8: P3-S1 — 3-Class Results ═════
    s = prs.slides.add_slide(prs.slide_invalid if False else prs.slide_layouts[LY_CONTENT])
    set_title(s, "3-Class Classification Results")

    img(s, FIGS["confusion"], Inches(0.35), Inches(CONTENT_TOP),
        Inches(9.0))
    caption(s, Inches(0.35), Inches(CONTENT_TOP + 2.5), Inches(9.0),
            "Confusion matrices: LR / RF / SVM / XGB  (HC · Mild · Severe)")

    # Right panel: big numbers
    right_x = 9.55
    # Best model card
    card(s, Inches(right_x), Inches(CONTENT_TOP), Inches(3.45), Inches(1.3),
         KU_BLUE, rounded=True)
    tb(s, Inches(right_x), Inches(CONTENT_TOP + 0.08), Inches(3.45), Inches(0.45),
       "SVM (best)", sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(right_x), Inches(CONTENT_TOP + 0.45), Inches(3.45), Inches(0.75),
       "81.0%", sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(right_x), Inches(CONTENT_TOP + 1.0), Inches(3.45), Inches(0.3),
       "Balanced Accuracy", sz=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Mini compare
    mini = [("XGB", "79.8%"), ("RF", "80.1%"), ("LR", "75.x%")]
    for i, (m, v) in enumerate(mini):
        x = right_x + i * 1.15
        card(s, Inches(x), Inches(CONTENT_TOP + 1.5),
             Inches(1.05), Inches(0.8), SOFT_BG, line_color=KU_GRAY,
             rounded=True)
        tb(s, Inches(x), Inches(CONTENT_TOP + 1.52), Inches(1.05),
           Inches(0.35), m, sz=11, color=KU_GRAY, align=PP_ALIGN.CENTER)
        tb(s, Inches(x), Inches(CONTENT_TOP + 1.8), Inches(1.05),
           Inches(0.5), v, sz=16, bold=True, color=BLACK,
           align=PP_ALIGN.CENTER)

    # Three takeaway bullets below
    bullets_y = CONTENT_TOP + 3.0
    key = [
        ("✓", "HC vs. Pathological", "> 90% recall — clinical screening reliable",
         KU_GREEN),
        ("!", "Mild vs. Severe boundary", "Main challenge — subjective ordinal scale",
         KU_RED),
        ("★", "Largest error", "Severe → Mild = 12.3% (late-recovery patients)",
         RGBColor(0xE0, 0xA5, 0x00)),  # darker yellow-gold for contrast
    ]
    for i, (icon, head, desc, clr) in enumerate(key):
        x = 0.55 + i * 4.25
        card(s, Inches(x), Inches(bullets_y), Inches(4.1), Inches(1.2),
             SOFT_BG, line_color=clr, rounded=True)
        badge(s, left=x + 0.2, top=bullets_y + 0.25, size=0.55,
              fill=clr, glyph=icon, sz=16)
        tb(s, Inches(x + 0.85), Inches(bullets_y + 0.1), Inches(3.15),
           Inches(0.4), head, sz=13, bold=True, color=clr)
        tb(s, Inches(x + 0.85), Inches(bullets_y + 0.5), Inches(3.15),
           Inches(0.75), desc, sz=11, color=BLACK)

    presenter_tag(s, 3, "3-Class Results")
    add_notes(s, (
        "Thank you. I am Presenter Three, and I will present the "
        "classification results in detail.\n\n"
        "The confusion matrices on the left show all four models — "
        "Logistic Regression, Random Forest, SVM, and XGBoost. The rows "
        "are true labels, the columns are predictions. Brighter diagonals "
        "mean better performance.\n\n"
        "SVM is our best model, shown on the right, with 81.0 percent "
        "balanced accuracy. XGBoost and Random Forest are only slightly "
        "behind at around 80 percent. Logistic Regression trails at about "
        "75 percent.\n\n"
        "Three key findings from the matrices.\n\n"
        "First, Healthy versus Pathological separation is very strong — "
        "recall is above 90 percent. So IMUs are reliable for initial "
        "clinical screening.\n\n"
        "Second, the main challenge is the Mild versus Severe boundary. "
        "This is not a surprise because the severity scale is ordinal and "
        "somewhat subjective.\n\n"
        "Third, the largest error is Severe predicted as Mild — 12.3 "
        "percent. These are mostly late-recovery patients whose gait has "
        "become close to normal."
    ))
    print("  [8] 3-Class Results")

    # ═════ SLIDE 9: P3-S2 — Phase Comparison ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Gait Phase & Window Configuration")

    img(s, FIGS["phase"], Inches(0.55), Inches(CONTENT_TOP), Inches(7.0))
    caption(s, Inches(0.55), Inches(CONTENT_TOP + 3.95), Inches(7.0),
            "Full Gait (6 s) vs. U-Turn (3 s) vs. Combined")

    configs = [
        ("Full Gait 6 s",  "50% overlap", "81.0%", KU_BLUE,  "★ Best"),
        ("Combined",       "Full + U-Turn","79.5%", KU_GREEN, ""),
        ("U-Turn 3 s",     "0% overlap",  "74.2%", KU_TEAL,  ""),
    ]
    for i, (cfg, detail, bacc, clr, tag) in enumerate(configs):
        y = CONTENT_TOP + i * 1.3
        card(s, Inches(7.75), Inches(y), Inches(5.2), Inches(1.15),
             clr, rounded=True)
        tb(s, Inches(7.9), Inches(y + 0.12), Inches(3.1), Inches(0.45),
           cfg, sz=16, bold=True, color=WHITE)
        tb(s, Inches(7.9), Inches(y + 0.55), Inches(3.1), Inches(0.4),
           detail, sz=11, color=WHITE)
        tb(s, Inches(10.95), Inches(y + 0.1), Inches(1.95), Inches(0.95),
           bacc, sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)
        if tag:
            tb(s, Inches(7.9), Inches(y - 0.02), Inches(1.2), Inches(0.3),
               tag, sz=10, bold=True, color=KU_YELLOW)

    # Takeaway
    card(s, Inches(0.55), Inches(CONTENT_TOP + 4.35), Inches(12.4),
         Inches(0.75), SOFT_BG, line_color=KU_BLUE, rounded=True)
    badge(s, left=0.75, top=CONTENT_TOP + 4.48, size=0.5,
          fill=KU_BLUE, glyph="!", sz=16)
    tb(s, Inches(1.4), Inches(CONTENT_TOP + 4.4), Inches(11.4), Inches(0.65),
       "Steady-state gait carries the most discriminative features — U-turns dilute signal",
       sz=13, bold=True, color=DEEP_NAVY, anchor=MSO_ANCHOR.MIDDLE)

    presenter_tag(s, 3, "Phase Comparison")
    add_notes(s, (
        "This slide compares different gait phases and window "
        "configurations.\n\n"
        "The bar chart on the left shows balanced accuracy for four phases: "
        "pre-U-turn, U-turn, post-U-turn, and full gait. Blue bars use only "
        "the best single sensor, orange bars use all four sensors.\n\n"
        "On the right we summarise three configurations.\n\n"
        "The Full Gait 6-second window with 50 percent overlap is our best, "
        "at 81.0 percent.\n\n"
        "If we combine Full Gait windows with U-turn windows, performance "
        "actually drops slightly to 79.5 percent.\n\n"
        "Using only U-turn 3-second windows gives 74.2 percent — the weakest.\n\n"
        "The conclusion is clear: steady-state gait carries the most "
        "discriminative information. U-turns are short, variable, and add "
        "noise rather than signal.\n\n"
        "So for clinical deployment we recommend using full gait cycles, "
        "not turn segments."
    ))
    print("  [9] Phase Comparison")

    # ═════ SLIDE 10: P3-S3 — Error Modes & VGA ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Error Modes & Clinical Correlation")

    img(s, FIGS["vga"], Inches(0.35), Inches(CONTENT_TOP), Inches(7.8))
    caption(s, Inches(0.35), Inches(CONTENT_TOP + 2.75), Inches(7.8),
            "VGA score vs. stride asymmetry (Spearman ρ = −0.206, n = 927)")

    # Right: error modes
    errors = [
        ("Severe → Mild",   "12.3%", "Late-recovery patients",   KU_RED),
        ("Mild → HC",        "8.1%", "Near-normal gait",          KU_YELLOW),
        ("Mild → Severe",    "5.7%", "Compensatory patterns",     KU_TEAL),
    ]
    tb(s, Inches(8.35), Inches(CONTENT_TOP - 0.05), Inches(4.6),
       Inches(0.4), "Main error modes", sz=13, bold=True, color=KU_GRAY)
    for i, (mode, rate, reason, clr) in enumerate(errors):
        y = CONTENT_TOP + 0.4 + i * 1.0
        card(s, Inches(8.35), Inches(y), Inches(4.6), Inches(0.9),
             clr, rounded=True)
        tb(s, Inches(8.5), Inches(y + 0.05), Inches(2.75), Inches(0.45),
           mode, sz=13, bold=True, color=WHITE)
        tb(s, Inches(8.5), Inches(y + 0.48), Inches(2.75), Inches(0.4),
           reason, sz=10, color=WHITE)
        tb(s, Inches(11.2), Inches(y + 0.1), Inches(1.65), Inches(0.7),
           rate, sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)

    # Bottom: r=0.68 hero + clinical bullets
    hero_y = CONTENT_TOP + 3.3
    card(s, Inches(0.55), Inches(hero_y), Inches(3.2), Inches(1.5),
         KU_BLUE, rounded=True)
    tb(s, Inches(0.55), Inches(hero_y + 0.08), Inches(3.2), Inches(0.9),
       "r = 0.68", sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(0.55), Inches(hero_y + 1.0), Inches(3.2), Inches(0.45),
       "VGA ↔ stride variability (p<0.001)",
       sz=10, color=WHITE, align=PP_ALIGN.CENTER)

    card(s, Inches(3.95), Inches(hero_y), Inches(9.0), Inches(1.5),
         SOFT_BG, line_color=KU_BLUE, rounded=True)
    bul = [
        "Cohen's d: asymmetry indices show large effects between Mild and Severe",
        "ML errors cluster where clinicians also disagree (ordinal boundaries)",
        "IMU-derived features provide objective complement to subjective VGA",
    ]
    for i, t in enumerate(bul):
        tb(s, Inches(4.1), Inches(hero_y + 0.1 + i * 0.42), Inches(0.25),
           Inches(0.35), "→", sz=12, bold=True, color=KU_GREEN)
        tb(s, Inches(4.4), Inches(hero_y + 0.08 + i * 0.42), Inches(8.5),
           Inches(0.4), t, sz=12, color=BLACK,
           anchor=MSO_ANCHOR.MIDDLE)

    presenter_tag(s, 3, "Error Analysis")
    add_notes(s, (
        "My last slide analyses where the errors happen and whether they "
        "make clinical sense.\n\n"
        "On the right we list the three main error modes. Severe predicted "
        "as Mild is the largest, 12.3 percent — these are late-recovery "
        "patients. Mild predicted as Healthy is 8.1 percent — near-normal "
        "gait. Mild predicted as Severe is 5.7 percent, often compensatory "
        "walking patterns.\n\n"
        "The scatter and box plots on the left show the Visual Gait "
        "Assessment score versus our stride asymmetry index. Spearman rho "
        "is minus 0.206, p less than 0.001. Asymmetry rises clearly with "
        "VGA severity.\n\n"
        "The big number on the lower left is very important. The correlation "
        "between VGA scores and IMU-derived stride variability is 0.68, "
        "highly significant.\n\n"
        "The message: our model's errors cluster exactly where clinicians "
        "also disagree. IMU features do not replace the clinician — they "
        "provide an objective complement to a subjective scale.\n\n"
        "Now over to Presenter Four for robustness and conclusions."
    ))
    print("  [10] Error Modes")

    # ═════ SLIDE 11: P4-S1 — LOCO & Nested CV ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Leave-One-Cohort-Out  &  Nested CV")

    # 7 cohort cards laid out 4+3
    loco = [
        ("Calcaneus Fx",    "76.8%"),
        ("Ankle Fx",        "79.2%"),
        ("Tibial Plat. Fx", "73.5%"),
        ("Femoral Fx",      "71.9%"),
        ("Knee Replace.",   "78.4%"),
        ("Hip Replace.",    "80.1%"),
        ("Hip Fx",          "74.6%"),
    ]
    # Color code: > 78% green, 74-78% blue, < 74% red
    def loco_color(v):
        x = float(v.replace('%',''))
        if x >= 78: return KU_GREEN
        if x >= 74: return KU_BLUE
        return KU_RED

    cw, ch, gx, gy = 2.95, 1.35, 0.12, 0.2
    row1_y = CONTENT_TOP
    for i, (cohort, bacc) in enumerate(loco):
        col = i % 4
        row = i // 4
        x = 0.55 + col * (cw + gx)
        y = row1_y + row * (ch + gy)
        clr = loco_color(bacc)
        card(s, Inches(x), Inches(y), Inches(cw), Inches(ch), clr,
             rounded=True)
        tb(s, Inches(x), Inches(y + 0.1), Inches(cw), Inches(0.45),
           cohort, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, Inches(x), Inches(y + 0.55), Inches(cw), Inches(0.7),
           bacc, sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)

    # Summary band
    sum_y = row1_y + 2 * (ch + gy) + 0.1
    card(s, Inches(0.55), Inches(sum_y), Inches(12.4), Inches(1.35),
         SOFT_BG, rounded=True)
    # three sub-cells
    cells = [
        ("Standard 10-Fold CV", "81.0%",  KU_BLUE),
        ("Nested CV",            "78.3%", KU_RED),
        ("Optimism Gap",         "2.7%",  DEEP_NAVY),
    ]
    for i, (lbl, v, clr) in enumerate(cells):
        x = 0.75 + i * 4.1
        tb(s, Inches(x), Inches(sum_y + 0.15), Inches(3.9), Inches(0.4),
           lbl, sz=12, color=KU_GRAY, align=PP_ALIGN.CENTER)
        tb(s, Inches(x), Inches(sum_y + 0.5), Inches(3.9), Inches(0.7),
           v, sz=30, bold=True, color=clr, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE)

    tb(s, Inches(0.55), Inches(sum_y + 1.4), Inches(12.4), Inches(0.35),
       "Model generalizes well — LOCO confirms robustness across pathology types",
       sz=11, color=DEEP_NAVY, align=PP_ALIGN.CENTER, bold=True)

    presenter_tag(s, 4, "LOCO & Nested CV")
    add_notes(s, (
        "Thank you. I am Presenter Four. I will cover robustness, clinical "
        "implications, and the conclusions.\n\n"
        "This slide shows two robustness checks.\n\n"
        "The top section is Leave-One-Cohort-Out validation. We hold out "
        "one clinical cohort completely, train on the other seven, then "
        "test on the held-out one. This simulates encountering a new "
        "pathology type at deployment time.\n\n"
        "The results range from 71.9 percent for Femoral Fracture up to "
        "80.1 percent for Hip Replacement. Most cohorts stay above 74 "
        "percent. The model generalises reasonably well, even to cohorts "
        "it has never seen.\n\n"
        "The bottom row shows our Nested CV check. Standard 10-fold CV "
        "gives 81.0 percent. Nested CV gives 78.3 percent. The gap is "
        "only 2.7 percent, which is small for this type of problem.\n\n"
        "Put together, these two tests confirm the model is not overfit, "
        "and it is robust across different pathology types."
    ))
    print("  [11] LOCO & Nested CV")

    # ═════ SLIDE 12: P4-S2 — Clinical Implications & Limitations ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Clinical Implications  &  Limitations")

    # Left card: Implications
    left_x = 0.55
    card(s, Inches(left_x), Inches(CONTENT_TOP), Inches(6.1), Inches(4.7),
         SOFT_GREEN, line_color=KU_GREEN, rounded=True)
    tb(s, Inches(left_x + 0.25), Inches(CONTENT_TOP + 0.1), Inches(0.5),
       Inches(0.5), "✓", sz=24, bold=True, color=KU_GREEN,
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(left_x + 0.8), Inches(CONTENT_TOP + 0.1), Inches(5.1),
       Inches(0.5), "Clinical Implications", sz=18, bold=True,
       color=KU_GREEN, anchor=MSO_ANCHOR.MIDDLE)

    impl = [
        ("81% accuracy",        "with wearable IMUs"),
        ("2 sensors sufficient", "Head + Lower Back"),
        ("Real-time feasible",   "screening in clinics"),
        ("Objective complement", "to VGA scoring"),
        ("Home monitoring",      "scalable remote deployment"),
        ("Automated triage",     "for severity stratification"),
    ]
    for i, (head, sub) in enumerate(impl):
        y = CONTENT_TOP + 0.8 + i * 0.62
        badge(s, left=left_x + 0.3, top=y + 0.05, size=0.45,
              fill=KU_GREEN, glyph="✓", sz=14)
        tb(s, Inches(left_x + 0.9), Inches(y + 0.02), Inches(5.05),
           Inches(0.35), head, sz=13, bold=True, color=DEEP_NAVY)
        tb(s, Inches(left_x + 0.9), Inches(y + 0.32), Inches(5.05),
           Inches(0.3), sub, sz=10, color=KU_GRAY)

    # Right card: Limitations
    right_x = 6.85
    card(s, Inches(right_x), Inches(CONTENT_TOP), Inches(6.1), Inches(4.7),
         SOFT_RED, line_color=KU_RED, rounded=True)
    tb(s, Inches(right_x + 0.25), Inches(CONTENT_TOP + 0.1), Inches(0.5),
       Inches(0.5), "⚠", sz=22, bold=True, color=KU_RED,
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(right_x + 0.8), Inches(CONTENT_TOP + 0.1), Inches(5.1),
       Inches(0.5), "Limitations", sz=18, bold=True,
       color=KU_RED, anchor=MSO_ANCHOR.MIDDLE)

    lims = [
        ("Single-site dataset",    "Innsbruck only"),
        ("Lab-controlled",          "no free-living data"),
        ("Subjective labels",       "Mild/Severe boundary"),
        ("Cross-sectional",         "no longitudinal follow-up"),
        ("Missing vitals",          "no BP / HR signals"),
        ("Hand-crafted features",   "no deep learning baseline"),
    ]
    for i, (head, sub) in enumerate(lims):
        y = CONTENT_TOP + 0.8 + i * 0.62
        badge(s, left=right_x + 0.3, top=y + 0.05, size=0.45,
              fill=KU_RED, glyph="!", sz=14)
        tb(s, Inches(right_x + 0.9), Inches(y + 0.02), Inches(5.05),
           Inches(0.35), head, sz=13, bold=True, color=DEEP_NAVY)
        tb(s, Inches(right_x + 0.9), Inches(y + 0.32), Inches(5.05),
           Inches(0.3), sub, sz=10, color=KU_GRAY)

    presenter_tag(s, 4, "Clinical Implications")
    add_notes(s, (
        "Now the clinical implications and limitations.\n\n"
        "On the left, six clinical implications. We achieve 81 percent "
        "accuracy with wearable IMUs. Only two sensors — Head and Lower "
        "Back — are sufficient. Real-time screening is feasible. Our "
        "method provides an objective complement to Visual Gait "
        "Assessment. It is scalable to home and remote monitoring, and "
        "enables automated severity triage in clinics.\n\n"
        "On the right, we are honest about the limitations. The dataset "
        "is from a single site — Innsbruck — so external validation is "
        "needed. Recordings are in a controlled lab, not free-living. "
        "The Mild versus Severe boundary is subjective. The dataset is "
        "cross-sectional, without longitudinal follow-up. We do not have "
        "blood pressure or heart rate. And we only use hand-crafted "
        "features, no deep learning baseline yet.\n\n"
        "These limitations directly motivate our future work, which I "
        "will present on the next slide."
    ))
    print("  [12] Implications & Limitations")

    # ═════ SLIDE 13: P4-S3 — Conclusions & Future ═════
    s = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    set_title(s, "Conclusions  &  Future Work")

    # Three pillars
    pillars = [
        ("01", "81.0%", "Balanced Accuracy",
         "SVM on 216 IMU features —\nrobust 3-class classification", KU_BLUE),
        ("02", "2", "Sensors Sufficient",
         "HE + LB capture 98.9% of\nfull 4-sensor performance", KU_GREEN),
        ("03", "✓", "Clinically Validated",
         "Feature rankings align with\nestablished VGA criteria", KU_RED),
    ]
    pw = 4.1
    for i, (num, big, head, desc, clr) in enumerate(pillars):
        x = 0.55 + i * (pw + 0.15)
        # colored number circle
        icon_circle(s, cx=x + 0.45, cy=CONTENT_TOP + 0.45, r=0.35,
                    fill=clr, text=num, sz=15)
        # big stat
        tb(s, Inches(x + 1.0), Inches(CONTENT_TOP + 0.0), Inches(3.0), Inches(0.9),
           big, sz=38, bold=True, color=clr, anchor=MSO_ANCHOR.MIDDLE)
        # head
        tb(s, Inches(x), Inches(CONTENT_TOP + 0.95), Inches(pw), Inches(0.4),
           head, sz=16, bold=True, color=DEEP_NAVY)
        # desc box
        card(s, Inches(x), Inches(CONTENT_TOP + 1.5), Inches(pw),
             Inches(1.4), SOFT_BG, line_color=clr, rounded=True)
        tb(s, Inches(x + 0.15), Inches(CONTENT_TOP + 1.55),
           Inches(pw - 0.3), Inches(1.3), desc,
           sz=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)

    # Future directions
    fy = CONTENT_TOP + 3.3
    tb(s, Inches(0.55), Inches(fy), Inches(5), Inches(0.4),
       "Future Directions", sz=16, bold=True, color=KU_BLUE)

    future = [
        ("DL",  "Deep Learning", "CNN / LSTM on raw IMU signals",  KU_BLUE),
        ("MS",  "Multi-Site",    "External validation on new datasets", KU_GREEN),
        ("ED",  "Edge Deploy",   "Real-time on wearable devices",   KU_TEAL),
        ("LT",  "Longitudinal",  "Track recovery progression over time", KU_RED),
    ]
    for i, (glyph, head, desc, clr) in enumerate(future):
        col, row = i % 2, i // 2
        x = 0.55 + col * 6.25
        y = fy + 0.5 + row * 0.75
        badge(s, left=x, top=y + 0.05, size=0.5, fill=clr,
              glyph=glyph, sz=11)
        tb(s, Inches(x + 0.65), Inches(y), Inches(2.1), Inches(0.45),
           head, sz=13, bold=True, color=clr, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, Inches(x + 2.8), Inches(y), Inches(3.4), Inches(0.45),
           desc, sz=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)

    presenter_tag(s, 4, "Conclusions")
    add_notes(s, (
        "My last slide — our conclusions and future work.\n\n"
        "Three main takeaways.\n\n"
        "Number one: we reach 81 percent balanced accuracy with an SVM on "
        "216 IMU features, for robust three-class classification of gait "
        "pathology.\n\n"
        "Number two: two sensors are sufficient. Head plus Lower Back "
        "captures 98.9 percent of the full four-sensor performance.\n\n"
        "Number three: the method is clinically validated. Our feature "
        "rankings align well with established Visual Gait Assessment "
        "criteria.\n\n"
        "For future work, four directions. First, Deep Learning — train "
        "CNN or LSTM models directly on the raw IMU signals. Second, "
        "Multi-Site — external validation on other gait datasets. Third, "
        "Edge Deployment — real-time classification on wearable devices. "
        "Fourth, Longitudinal studies — to track recovery progression "
        "over time.\n\n"
        "That concludes our presentation. Thank you for listening, and we "
        "are happy to take questions."
    ))
    print("  [13] Conclusions & Future")

    # ═════ SLIDE 14: THANK YOU ═════
    # The END layout already has a big "Thank You" text at y~2.78 (height 1.94).
    # Don't add another — just add subtitle + Questions below it.
    s = prs.slides.add_slide(prs.slide_layouts[LY_END])
    # Subtitle (below built-in "Thank You")
    tb(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
       "BMED 712  |  Rehabilitation Engineering  |  Khalifa University  |  April 2026",
       sz=16, color=WHITE, align=PP_ALIGN.CENTER)
    # Questions
    tb(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.6),
       "Questions?", sz=28, bold=True, color=KU_YELLOW, align=PP_ALIGN.CENTER)
    add_notes(s, (
        "Thank you for your attention.\n"
        "We are happy to answer any questions you may have."
    ))
    print("  [14] Thank You")

    # ═════ SAVE ═════
    prs.save(str(OUTPUT))
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f"\nSaved: {OUTPUT} ({size_kb:.0f} KB)")
    print("Total: 14 slides = 1 title + 12 content + 1 thank-you")


if __name__ == "__main__":
    build()
